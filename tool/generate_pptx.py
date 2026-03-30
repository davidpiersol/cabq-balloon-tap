#!/usr/bin/env python3
"""Generate the Balloon Tap 2.0 concept deck (PowerPoint)."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DOCS = Path(__file__).resolve().parent.parent / "docs"
CONCEPTS = DOCS / "concepts"
OUT = DOCS / "Balloon_Tap_2.0_Concepts.pptx"

SLIDES = [
    {
        "title": "Balloon Tap 2.0",
        "subtitle": "Concept design review — mass ascension title (frame 1 references)",
        "image": None,
    },
    {
        "title": "1 — Mass Ascension Dawn (Splash)",
        "subtitle": "Pre-dawn NM sky, mass ascension balloons, hero flame, title CTA.",
        "image": CONCEPTS / "v2_mockup_1_mass_ascension_dawn.png",
    },
    {
        "title": "2 — In-Flight Gameplay",
        "subtitle": "Active play: burner flame, glass HUD, desert parallax, hint text.",
        "image": CONCEPTS / "v2_mockup_2_inflight_gameplay.png",
    },
    {
        "title": "3 — Sandia Sunset Skin",
        "subtitle": "Orange/magenta/purple sky, sunset balloon, red chile collectible.",
        "image": CONCEPTS / "v2_mockup_3_sandia_sunset.png",
    },
    {
        "title": "4 — Game Over / New Best",
        "subtitle": "Celebration card, gold confetti, score typography, Play again CTA.",
        "image": CONCEPTS / "v2_mockup_4_gameover_newbest.png",
    },
    {
        "title": "5 — First-Run Onboarding",
        "subtitle": "Hold/release/coast instructions, Balloon Fiesta backdrop, Let's fly.",
        "image": CONCEPTS / "v2_mockup_5_onboarding.png",
    },
    {
        "title": "Design Palette",
        "subtitle": (
            "Sky Top #87CEEB  |  Horizon #FFE8CC  |  Sky Bottom #E8B86D\n"
            "Sandia Pink #FF8FA3  |  Sunset Orange #FFB347  |  Magenta #6B2D5C\n"
            "Primary #0A5F73  |  Accent #D94E1F  |  Sand #F4EDE4\n"
            "Chile Red #B71C1C  |  Chile Green #2E7D32  |  Gold #FFD700"
        ),
        "image": None,
    },
    {
        "title": "Next Steps",
        "subtitle": (
            "- Replace placeholder Lottie/Rive with original NM artwork\n"
            "- Build and test on iOS Simulator + device\n"
            "- v3: non-intrusive sponsor banner ads (reserved layout)\n"
            "- Store screenshots from concept frames\n"
            "- Push to balloon-tap-2 GitHub repo"
        ),
        "image": None,
    },
]

CABQ_TEAL = RGBColor(0x0A, 0x5F, 0x73)
CABQ_ACCENT = RGBColor(0xD9, 0x4E, 0x1F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x2E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment

    for line in text.split("\n")[1:]:
        p2 = tf.add_paragraph()
        p2.text = line
        p2.font.size = Pt(font_size)
        p2.font.bold = False
        p2.font.color.rgb = color
        p2.alignment = alignment
    # fix: first paragraph only has first line
    p.text = text.split("\n")[0]


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # blank

    for s in SLIDES:
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, DARK)

        # Title
        add_text(slide, s["title"],
                 Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
                 font_size=36, bold=True, color=CABQ_ACCENT, alignment=PP_ALIGN.LEFT)

        if s["image"] and s["image"].exists():
            img_left = Inches(0.6)
            img_top = Inches(1.5)
            img_height = Inches(5.2)
            slide.shapes.add_picture(str(s["image"]), img_left, img_top, height=img_height)

            add_text(slide, s["subtitle"],
                     Inches(7.5), Inches(1.8), Inches(5.2), Inches(4),
                     font_size=20, color=WHITE, alignment=PP_ALIGN.LEFT)
        else:
            add_text(slide, s["subtitle"],
                     Inches(0.6), Inches(1.6), Inches(12), Inches(5),
                     font_size=24, color=WHITE, alignment=PP_ALIGN.LEFT)

    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    main()
